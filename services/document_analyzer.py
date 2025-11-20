import os
import tempfile

import olefile
import requests
from PyPDF2 import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from services.llm import refine


def download_file(url: str) -> str:
    """
    Presigned URL에서 파일 다운로드
    """
    response = requests.get(url, timeout=60)
    response.raise_for_status()

    # 임시 파일로 저장
    with tempfile.NamedTemporaryFile(delete=False, suffix='.tmp') as tmp_file:
        tmp_file.write(response.content)
        return tmp_file.name


def extract_text_from_pdf(file_path: str) -> str:
    """
    PDF에서 텍스트 추출
    """
    reader = PdfReader(file_path)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return text.strip()


def extract_text_from_docx(file_path: str) -> str:
    """
    DOCX에서 텍스트 추출
    """
    doc = Document(file_path)
    text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
    return text.strip()


def extract_text_from_xlsx(file_path: str) -> str:
    """
    XLSX에서 텍스트 추출 (첫 시트의 데이터)
    """
    wb = load_workbook(file_path, read_only=True)
    ws = wb.active

    text_lines = []
    for row in ws.iter_rows(values_only=True):
        row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
        if row_text.strip():
            text_lines.append(row_text)

    return "\n".join(text_lines)


def extract_text_from_pptx(file_path: str) -> str:
    """
    PPTX에서 텍스트 추출
    """
    prs = Presentation(file_path)
    text_lines = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_lines.append(shape.text)

    return "\n".join(text_lines).strip()


def extract_text_from_hwp(file_path: str) -> str:
    """
    HWP에서 텍스트 추출 (간단한 방법)
    """
    try:
        if olefile.isOleFile(file_path):
            ole = olefile.OleFileIO(file_path)
            # HWP 파일 구조에서 텍스트 추출 시도
            # 주의: 완벽하지 않을 수 있음
            text_streams = []
            for stream in ole.listdir():
                stream_name = "/".join(stream)
                if "BodyText" in stream_name or "PrvText" in stream_name:
                    try:
                        data = ole.openstream(stream).read()
                        # UTF-16으로 디코딩 시도
                        text = data.decode('utf-16', errors='ignore')
                        text_streams.append(text)
                    except:
                        pass
            ole.close()
            return "\n".join(text_streams).strip()
        else:
            # HWPX (zip 기반)는 복잡하므로 기본 메시지 반환
            return "[HWP 파일: 텍스트 추출 제한적]"
    except Exception as e:
        return f"[HWP 파일 분석 실패: {str(e)}]"


def extract_text_from_txt(file_path: str) -> str:
    """
    TXT 파일에서 텍스트 추출
    """
    encodings = ['utf-8', 'cp949', 'euc-kr', 'latin-1']

    for encoding in encodings:
        try:
            with open(file_path, 'r', encoding=encoding) as f:
                return f.read()
        except UnicodeDecodeError:
            continue

    # 모든 인코딩 실패 시
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()


def extract_text(file_path: str, extension: str) -> str:
    """
    확장자에 따라 텍스트 추출
    """
    ext = extension.lower()

    if ext == 'pdf':
        return extract_text_from_pdf(file_path)
    elif ext in ['doc', 'docx']:
        return extract_text_from_docx(file_path)
    elif ext in ['xls', 'xlsx']:
        return extract_text_from_xlsx(file_path)
    elif ext in ['ppt', 'pptx']:
        return extract_text_from_pptx(file_path)
    elif ext in ['hwp', 'hwpx']:
        return extract_text_from_hwp(file_path)
    elif ext in ['txt', 'log', 'md', 'markdown']:
        return extract_text_from_txt(file_path)
    elif ext in ['json', 'xml', 'yaml', 'yml', 'csv']:
        return extract_text_from_txt(file_path)
    else:
        raise ValueError(f"Unsupported file extension: {extension}")


def summarize_document(text: str, language: str, file_name: str) -> str:
    """
    LLM으로 문서 요약 생성
    """
    # 텍스트가 너무 길면 앞부분만 사용 (LLM 입력 제한)
    max_length = 10000
    if len(text) > max_length:
        text = text[:max_length] + "\n\n[문서가 잘렸습니다...]"

    # 언어 매핑
    lang_map = {
        'KO': 'Korean',
        'EN': 'English',
        'JA': 'Japanese',
        'FR': 'French'
    }
    target_language = lang_map.get(language.upper(), 'English')

    # LLM 프롬프트
    prompt = f"""Analyze and summarize this document in {target_language}.

File: {file_name}

Content:
{text}

Please provide:
1. Main topic (1-2 sentences)
2. Key points (3-5 bullet points)
3. Important keywords

Format the summary to be clear and concise for a developer chat context."""

    # LLM 호출 (refine 함수 재사용)
    summary = refine(text, text, target_language.lower())

    return summary


def analyze_document(file_url: str, file_name: str, summary_language: str) -> str:
    """
    문서 분석 메인 함수
    
    Args:
        file_url: Presigned GET URL
        file_name: 원본 파일명
        summary_language: 요약 언어 (KO, EN, JA, FR)
    
    Returns:
        요약 텍스트
    """
    temp_file = None
    try:
        # 1. 파일 다운로드
        temp_file = download_file(file_url)

        # 2. 확장자 추출
        extension = file_name.split('.')[-1] if '.' in file_name else ''

        # 3. 텍스트 추출
        text = extract_text(temp_file, extension)

        if not text or len(text.strip()) < 10:
            return f"[{file_name}: 텍스트 추출 실패 또는 내용 없음]"

        # 4. LLM 요약 생성
        summary = summarize_document(text, summary_language, file_name)

        return summary

    finally:
        # 임시 파일 삭제
        if temp_file and os.path.exists(temp_file):
            try:
                os.unlink(temp_file)
            except:
                pass